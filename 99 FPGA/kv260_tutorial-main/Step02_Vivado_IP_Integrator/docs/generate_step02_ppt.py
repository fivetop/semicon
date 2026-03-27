#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = "/home/woody/SoGang/2026-03/kv260_vivado_tutorial/Step02_Vivado_IP_Integrator/docs"

PRIMARY = RGBColor(0, 51, 102)
ACCENT = RGBColor(255, 102, 0)
LIGHT = RGBColor(240, 248, 255)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(51, 51, 51)
NAVY_DARK = RGBColor(0, 26, 64)
NAVY = RGBColor(0, 51, 102)
LIGHT_BLUE = RGBColor(170, 204, 238)

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY_DARK
    bg.line.fill.background()
    
    navy_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.5))
    navy_bar.fill.solid()
    navy_bar.fill.fore_color.rgb = NAVY
    navy_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12), Inches(1.5))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12), Inches(0.8))
    tb = sub_box.text_frame
    p = tb.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
    tb = footer.text_frame
    p = tb.paragraphs[0]
    p.text = "KV260 FPGA Training Series"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_BLUE

def section_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    number_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(2))
    tb = number_box.text_frame
    p = tb.paragraphs[0]
    p.text = number
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.7), Inches(3), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12), Inches(2))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

def content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(24)
        p.space_before = Pt(18)

def two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(5.8), Inches(5))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = PRIMARY
    left_card.line.width = Pt(2)
    
    ltb = left_card.text_frame
    p = ltb.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    for item in left_items:
        p = ltb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.space_before = Pt(8)
    
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.8), Inches(5.8), Inches(5))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = ACCENT
    right_card.line.width = Pt(2)
    
    rtb = right_card.text_frame
    p = rtb.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    for item in right_items:
        p = rtb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.space_before = Pt(8)

def diagram_slide(prs, title, diagram_text):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    diagram_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(12.333), Inches(4.5))
    diagram_box.fill.solid()
    diagram_box.fill.fore_color.rgb = DARK
    diagram_box.line.fill.background()
    
    tb = diagram_box.text_frame
    tb.word_wrap = True
    
    p = tb.add_paragraph()
    p.text = diagram_text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 255, 0)
    p.font.name = "Courier New"
    p.space_before = Pt(10)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    title_slide(prs, "Step 02: Vivado IP Integrator", "PS-PL 인터페이스와 블록 자동화")
    
    section_slide(prs, "01", "학습 목표")
    content_slide(prs, "본 튜토리얼의 목표", [
        "Vivado IP Integrator 기본 사용법 습득",
        "PS-PL 인터페이스 이해",
        "블록 자동화 (Block Automation) 활용",
        "플랫폼 익스포트 (XSA) 이해"
    ])
    
    two_column_slide(prs, "주요 학습 개념",
        "PS/PL 구성", [
            "PS: Processing System (ARM 프로세서)",
            "PL: Programmable Logic (FPGA 논리)",
            "Block Automation: 자동 연결 기능",
            "Board Preset: 보드 최적화 설정"
        ],
        "프로젝트 유형", [
            "RTL Project: PL만 설계",
            "Embedded Project: PS+PL 설계",
            "KV260에서는 Embedded 또는 RTL with Zynq",
            "Extensible Vitis Platform 옵션"
        ])
    
    section_slide(prs, "02", "프로젝트 생성")
    content_slide(prs, "Vivado 프로젝트 생성 방법", [
        "1. source /tools/Xilinx/Vivado/2021.2/settings64.sh",
        "2. vivado &",
        "3. Create Project -> RTL Project",
        "4. 프로젝트 이름: kv260_ipi_demo",
        "5. 보드 선택: Kria KV260 Vision AI Starter Kit"
    ])
    
    section_slide(prs, "03", "Zynq UltraScale+ MPSoC")
    diagram_slide(prs, "PS 아키텍처", """+------------------------------------------+
|         Processing System (PS)          |
|  +----------+  +----------+  +--------+ |
|  |  ARM     |  |  ARM     |  |  RPU   | |
|  |Cortex-A53|  |Cortex-R5 |  |        | |
|  +----------+  +----------+  +--------+ |
|  +----------+  +------------------+     |
|  |  DDR4    |  |  GPU/VDMA       |     |
|  +----------+  +------------------+     |
+------------------------------------------+
|         Programmable Logic (PL)         |
+------------------------------------------+""")
    
    content_slide(prs, "클럭 및 리셋 설정", [
        "pl_clk0: 100MHz (주 클럭)",
        "pl_clk1 ~ pl_clk3: 추가 클럭 (선택)",
        "pl_resetn0: 주 리셋 신호",
        "DDR4: 2GB, QSPI: 부트 플래시"
    ])
    
    section_slide(prs, "04", "블록 자동화")
    content_slide(prs, "Run Block Automation이란?", [
        "Vivado가 PS IP의 연결을 자동으로 설정",
        "AXI 버스 자동 연결",
        "클럭 (pl_clk0) 자동 연결",
        "리셋 (pl_resetn0) 자동 연결",
        "KV260에서는 Board Preset 권장"
    ])
    
    two_column_slide(prs, "연결 방식 비교",
        "자동 연결", [
            "Run Connection Automation 사용",
            "Zynq PS --> GPIO 자동 연결",
            "빠르고简便",
            "기본 설정으로 충분할 때"
        ],
        "수동 연결", [
            "AXI BUS 직접 연결",
            "세밀한 제어 가능",
            "복잡하지만 유연함",
            "고급 설정이 필요할 때"
        ])
    
    section_slide(prs, "05", "Clocking Wizard")
    diagram_slide(prs, "클럭 라우팅", """PS (pl_clk0 100MHz) --> Clocking Wizard --> PL Logic
                           |
                      clk_out1 (100MHz)
                      clk_out2 (200MHz)
                      clk_out3 (400MHz)
                      clk_out4 (50MHz)""")
    
    content_slide(prs, "Clocking Wizard 설정 방법", [
        "1. BD에서 '+' 버튼 클릭",
        "2. 'Clocking Wizard' 검색",
        "3. clk_wiz 추가 및 설정",
        "4. 출력 주파수: 100/200/400/50MHz"
    ])
    
    section_slide(prs, "06", "Processor System Reset")
    diagram_slide(prs, "리셋 연결", """PS (pl_resetn0) --> proc_sys_reset_0 --> PL Logic
                        |
                   dcm_locked
            (클럭 잠금까지 리셋 유지)""")
    
    content_slide(prs, "리셋 처리", [
        "비동기 리셋은 클럭 에지에서 동기화 필요",
        "Processor System Reset이 자동으로 처리",
        "dcm_locked: 클럭 안정화 신호"
    ])
    
    section_slide(prs, "07", "플랫폼 인터페이스")
    content_slide(prs, "Platform Setup 탭", [
        "Tools -> Platform Setup",
        "PFM.CLK: 클럭 인터페이스 설정",
        "PFM.IRQ: 인터럽트 설정",
        "Interrupt ID 0-31: PS to PL",
        "Interrupt ID 32-91: PL to PS"
    ])
    
    section_slide(prs, "08", "PS-PL 인터페이스")
    content_slide(prs, "주요 인터페이스", [
        "M_AXI_HPM0_FPD: 고성능 AXI (Full Power Domain)",
        "M_AXI_HPM0_LPD: 저전력 Domain AXI",
        "S_AXI_HPx_FPD: DDR 메모리 직접 접근 (0~3)",
        "pl_clk0, pl_resetn0: 기본 클럭/리셋"
    ])
    
    section_slide(prs, "09", "HDL Wrapper")
    two_column_slide(prs, "Wrapper 모드",
        "Vivado Managed", [
            "BD 변경 시 자동 재생성",
            "간단한 프로젝트에 적합",
            "Vivado가 관리"
        ],
        "User Managed", [
            "사용자가 직접 편집 가능",
            "고급 커스터마이징 가능",
            "수동 관리"
        ])
    
    section_slide(prs, "10", "하드웨어 익스포트")
    content_slide(prs, "XSA 파일 생성", [
        "방법 1: File -> Export -> Export Hardware",
        "방법 2: TCL: write_hw_platform -force -file kv260_platform.xsa",
        "Pre-synthesis: RTL 시뮬레이션용",
        "Post-synthesis: 비트스트림 생성용",
        "Vitis에서 플랫폼으로 사용: vitis --platform ~/kv260_platform.xsa"
    ])
    
    section_slide(prs, "11", "검증 체크리스트")
    content_slide(prs, "프로젝트 완료 조건", [
        "Zynq PS가 블록 디자인에 추가됨",
        "Block Automation이 정상 동작",
        "Clocking Wizard가 클럭 생성",
        "Processor System Reset이 연결됨",
        "HDL Wrapper가 생성됨",
        "XSA 파일이 익스포트됨"
    ])
    
    section_slide(prs, "12", "심화 연습 문제")
    content_slide(prs, "연습 문제", [
        "Clocking Wizard 설정: 다양한 출력 주파수 설정",
        "AXI 브릿지: PS와 PL 간 AXI 통신 설정",
        "인터럽트: PL에서 PS로 인터럽트送信 시스템 구성"
    ])
    
    title_slide(prs, "Step 02 완료", "다음 단계: AXI GPIO Integration")
    
    output_path = os.path.join(OUTPUT_DIR, "step02-ip-integrator.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")

if __name__ == "__main__":
    create_presentation()