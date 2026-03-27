#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = "/home/woody/SoGang/2026-03/kv260_vivado_tutorial/Step00_Overview/docs"

PRIMARY = RGBColor(0, 51, 102)
ACCENT = RGBColor(255, 102, 0)
LIGHT = RGBColor(240, 248, 255)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)
DARK = RGBColor(51, 51, 51)

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.8), Inches(13.333), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tb = sub_box.text_frame
    p = tb.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def section_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1.5))
    tb = num_box.text_frame
    p = tb.paragraphs[0]
    p.text = number
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(3), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

def content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(22)
        p.space_before = Pt(14)

def two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(5.8), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = PRIMARY
    left_card.line.width = Pt(2)
    
    ltb = left_card.text_frame
    p = ltb.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    for item in left_items:
        p = ltb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_before = Pt(10)
    
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5.8), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = ACCENT
    right_card.line.width = Pt(2)
    
    rtb = right_card.text_frame
    p = rtb.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    for item in right_items:
        p = rtb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_before = Pt(10)

def code_slide(prs, title, code_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
    code_box.line.fill.background()
    
    tb = code_box.text_frame
    tb.word_wrap = True
    
    for line in code_lines:
        p = tb.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 255, 0)
        p.font.name = "Courier New"

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    title_slide(prs, "VSCode Verilog 개발 환경", "가벼운 환경에서 Verilog 코딩 및 시뮬레이션")
    
    section_slide(prs, "01", "개요")
    content_slide(prs, "왜 VSCode인가?", [
        "Vivado 없이 가벼운 개발 환경",
        "문법 하이라이팅 및 자동 완성",
        "실시간 린팅 및 포맷팅",
        "시뮬레이션 지원 (Icarus Verilog + GTKWave)"
    ])
    
    section_slide(prs, "02", "VSCode 설치")
    content_slide(prs, "다운로드 위치", [
        "공식 웹사이트: code.visualstudio.com",
        "Download 버튼 클릭",
        "OS 선택 (Windows/Linux/Mac)",
        ".deb (Ubuntu), .rpm (Fedora), .dmg (Mac)"
    ])
    
    two_column_slide(prs, "설치 방법", "Ubuntu (APT 저장소)", [
        "wget -qO- https://packages.microsoft.com/keys/microsoft.asc | gpg --dearmor",
        "sudo install -D -o root -g root -m 644 microsoft.gpg /usr/share/keyrings/",
        'echo "deb [arch=amd64 signed-by=/usr/share/keyrings/microsoft.gpg] https://packages.microsoft.com/repos/code stable main" | sudo tee /etc/apt/sources.list.d/vscode.list',
        "sudo apt update && sudo apt install code"
    ], "Ubuntu (.deb 파일)", [
        "code.visualstudio.com에서 .deb 다운로드",
        "cd ~/Downloads",
        "sudo apt install ./code_*.deb",
        "code --version (확인)"
    ])
    
    two_column_slide(prs, "Windows/Mac 설치", "Windows", [
        "code.visualstudio.com에서 User Installer 다운로드",
        "실행 → 라이선스 동의 → 설치 위치 선택",
        "바탕화면에 바로가기 추가",
        "Code → Install 'code' command (PATH 추가)"
    ], "macOS", [
        "code.visualstudio.com에서 Universal 다운로드",
        "Download 폴더에서 .zip 압축 해제",
        "Visual Studio Code.app을 /Applications으로 이동",
        "Command Palette → Shell Command: Install 'code'"
    ])
    
    section_slide(prs, "03", "확장 프로그램")
    content_slide(prs, "필수 확장 프로그램", [
        "Verilog-HDL (mshr-h) - 문법 하이라이팅",
        "SystemVerilog (AndrewNolte) - linting",
        "Verilog Formatter - 코드 포맷팅"
    ])
    
    two_column_slide(prs, "확장 프로그램 설치", "VSCode에서", [
        "Ctrl+Shift+X로 확장 탭 열기",
        "검색창에 이름 입력",
        "Install 클릭"
    ], "Marketplace", [
        "Verilog-HDL: 130만+ 설치",
        "SystemVerilog: 무료",
        "Marketplace.visualstudio.com"
    ])
    
    section_slide(prs, "04", "유틸리티 설치")
    two_column_slide(prs, "필수 도구", "시뮬레이터", [
        "Icarus Verilog (iverilog)",
        "GTKWave ( waveform 뷰어)",
        "sudo apt install iverilog gtkwave"
    ], "symbol 분석", [
        "Universal Ctags",
        "정의 이동, 참조 찾기",
        "brew install universal-ctags (macOS)"
    ])
    
    section_slide(prs, "05", "VSCode 설정")
    code_slide(prs, "settings.json 예시", [
        '{',
        '  "verilog.ctags.path": "/usr/bin/ctags-universal",',
        '  "verilog.lint.slang.enabled": true,',
        '  "editor.formatOnSave": true,',
        '  "[verilog]": {',
        '    "editor.defaultFormatter": "bmpenuelas..."',
        '  }',
        '}'
    ])
    
    section_slide(prs, "06", "시뮬레이션")
    content_slide(prs, "실행 방법", [
        "iverilog -o design design.v tb_design.v",
        "vvp design (시뮬레이션 실행)",
        "gtkwave design.vcd (波形 확인)"
    ])
    
    code_slide(prs, "테스트bench 예시", [
        'module tb_dff;',
        '  reg d, clk, rstn;',
        '  wire q;',
        '  ',
        '  dff uut (.d(d), .clk(clk), .rstn(rstn), .q(q));',
        '  ',
        '  initial begin',
        '    clk = 0;',
        '    forever #5 clk = ~clk;',
        '  end',
        '  ',
        '  initial begin',
        '    rstn = 0; d = 0;',
        '    #10 rstn = 1;',
        '    #20 $finish;',
        '  end',
        'endmodule'
    ])
    
    section_slide(prs, "07", "단축키")
    content_slide(prs, "유용한 단축키", [
        "F12 - 정의로 이동",
        "Shift+F12 - 참조 찾기",
        "Ctrl+Shift+P - Command Palette",
        "Ctrl+Shift+P → Format Document - 포맷팅"
    ])
    
    section_slide(prs, "08", "프로젝트 구조")
    code_slide(prs, "권장 디렉토리 구조", [
        'project/',
        '├── src/         # RTL 소스',
        '│   ├── dff.v',
        '│   └── adder.v',
        '├── sim/         # 테스트벤치',
        '│   └── tb_dff.v',
        '├── constraints/ # XDC',
        '│   └── pin.xdc',
        '└── .vscode/    # 설정',
        '    └── settings.json'
    ])
    
    section_slide(prs, "09", "요약")
    content_slide(prs, "정리", [
        "VSCode + 확장 프로그램으로 가벼운 개발 환경",
        "Icarus Verilog + GTKWave로 시뮬레이션",
        "Universal Ctags로 symbol 분석",
        "Vivado 없이도 충분히 실습 가능"
    ])
    
    title_slide(prs, "다음 단계", "Vivado 프로젝트로 FPGA 합성하기")
    
    output_path = os.path.join(OUTPUT_DIR, "vscode-verilog.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")

if __name__ == "__main__":
    create_presentation()
