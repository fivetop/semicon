#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = "/home/woody/SoGang/2026-03/kv260_vivado_tutorial/Step00_Overview/docs"

PRIMARY = RGBColor(0, 51, 102)
ACCENT = RGBColor(255, 102, 0)
LIGHT = RGBColor(240, 248, 255)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)
DARK = RGBColor(51, 51, 51)

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.8), Inches(13.333), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tb = sub_box.text_frame
    p = tb.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def section_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1.5))
    tb = num_box.text_frame
    p = tb.paragraphs[0]
    p.text = number
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(3), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

def content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(22)
        p.space_before = Pt(14)

def two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(5.8), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = PRIMARY
    left_card.line.width = Pt(2)
    
    ltb = left_card.text_frame
    p = ltb.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    for item in left_items:
        p = ltb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_before = Pt(10)
    
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5.8), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = ACCENT
    right_card.line.width = Pt(2)
    
    rtb = right_card.text_frame
    p = rtb.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    for item in right_items:
        p = rtb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_before = Pt(10)

def code_slide(prs, title, code_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
    code_box.line.fill.background()
    
    tb = code_box.text_frame
    tb.word_wrap = True
    
    for line in code_lines:
        p = tb.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 255, 0)
        p.font.name = "Courier New"

def workflow_slide(prs, title, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    for i, step in enumerate(steps):
        x = Inches(0.5 + i * 2.5)
        y = Inches(2.2)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = PRIMARY
        card.line.fill.background()
        
        tb = card.text_frame
        p = tb.paragraphs[0]
        p.text = step
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        if i < len(steps) - 1:
            arrow = slide.shapes.add_textbox(x + Inches(2.2), y + Inches(0.7), Inches(0.3), Inches(0.5))
            atb = arrow.text_frame
            ap = atb.paragraphs[0]
            ap.text = "→"
            ap.font.size = Pt(24)
            ap.font.color.rgb = DARK
            ap.alignment = PP_ALIGN.CENTER

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    title_slide(prs, "Git 기초와 명령어", "버전 관리 시스템의 기초")
    
    section_slide(prs, "01", "Git이란?")
    content_slide(prs, "Git 이해", [
        "버전 관리 시스템 (Version Control System)",
        "파일 변경 이력 추적",
        "多人 협업 지원",
        "Linux 커널 제작자 Linus Torvalds가 개발"
    ])
    
    two_column_slide(prs, "왜 Git인가?", "장점", [
        "분산 버전 관리",
        "빠른 속도",
        "브랜치 지원",
        "오픈소스"
    ], "용도", [
        "코드 버전 관리",
        "多人 협업",
        "이력 추적",
        "백업"
    ])
    
    section_slide(prs, "02", "설치")
    two_column_slide(prs, "설치 방법", "Ubuntu/Linux", [
        "sudo apt update",
        "sudo apt install git",
        "git --version",
        "git config --global"
    ], "Windows", [
        "git-scm.com 다운로드",
        " installer 실행",
        "기본 옵션 유지",
        "git Bash에서 확인"
    ])
    
    section_slide(prs, "03", "기본 설정")
    code_slide(prs, "초기 설정 명령어", [
        'git config --global user.name "Your Name"',
        'git config --global user.email "your@email.com"',
        'git config --list',
        'git config --global core.editor code'
    ])
    
    section_slide(prs, "04", "기초 Workflow")
    workflow_slide(prs, "기본 작업 흐름", [
        "Working\nDirectory",
        "Staging\nArea",
        "Repository",
        "Remote"
    ])
    
    two_column_slide(prs, "핵심 명령어", "git init", [
        "새 저장소 생성",
        "git init",
        "ls -la (/.git 확인)",
        "touch README.md && git add ."
    ], "git clone", [
        "既存 저장소 복제",
        "git clone <URL>",
        "git clone <URL> my-folder",
        "git pull (업데이트)"
    ])
    
    section_slide(prs, "05", "추가 및 커밋")
    code_slide(prs, "파일 추가 및 커밋", [
        'git status                    # 상태 확인',
        'git add <file>               # 파일 추가',
        'git add .                    # 모든 변경파일 추가',
        'git commit -m "메시지"       # 커밋',
        'git commit -am "메시지"      # 추가+커밋 (tracked 파일)'
    ])
    
    two_column_slide(prs, "커밋 관리", "히스토리 확인", [
        "git log",
        "git log --oneline",
        "git log -n 3",
        "git log --graph"
    ], "이전 커밋으로 이동", [
        "git checkout <commit>",
        "git checkout main (돌아오기)",
        "git revert <commit>",
        "git reset --hard <commit>"
    ])
    
    section_slide(prs, "06", "브랜치")
    code_slide(prs, "브랜치 명령어", [
        'git branch              # 브랜치 목록',
        'git branch <name>       # 새 브랜치 생성',
        'git checkout <branch>  # 브랜치 전환',
        'git checkout -b <name> # 생성+전환',
        'git merge <branch>     # 병합',
        'git branch -d <branch> # 삭제'
    ])
    
    section_slide(prs, "07", "원격 저장소")
    code_slide(prs, "원격 명령어", [
        'git remote -v                    # 원격 저장소 확인',
        'git remote add origin <URL>     # 원격 추가',
        'git push origin main           # 푸시',
        'git pull origin main            # 풀',
        'git fetch origin               # 패치'
    ])
    
    two_column_slide(prs, "Push/Pull", "git push", [
        "git push origin main",
        "git push -u origin main (최초)",
        "git push origin <branch>",
        "git push --force (주의!)"
    ], "git pull", [
        "git pull origin main",
        "git pull --rebase origin main",
        "git fetch + git merge",
        "conflict 해결 방법 배움"
    ])
    
    section_slide(prs, "08", "실용 명령어")
    code_slide(prs, "자주 사용하는 명령어", [
        'git status              # 상태',
        'git diff                # 변경사항',
        'git diff --staged       # 스테이징 영역',
        'git log --oneline -n 5 # 최근 5개',
        'git stash               # 임시 저장',
        'git stash pop          # 복원'
    ])
    
    two_column_slide(prs, "Undo", "작업 취소", [
        "git checkout -- <file>",
        "git restore <file>",
        "git restore --staged <file>",
        "git reset HEAD <file>"
    ], "커밋 취소", [
        "git reset --soft HEAD~1",
        "git reset --mixed HEAD~1",
        "git reset --hard HEAD~1 (주의!)",
        "git revert <commit>"
    ])
    
    section_slide(prs, "09", "협업 Workflow")
    content_slide(prs, "권장 협업 방식", [
        "main 브랜치는保护的",
        "기능 개발은 feature 브랜치에서",
        "Pull Request로 병합 요청",
        "코드 리뷰 후 merge"
    ])
    
    code_slide(prs, "기능 개발流程", [
        'git checkout -b feature/login',
        'working...',
        'git add . && git commit -m "Add login"',
        'git push origin feature/login',
        'Pull Request 생성 → 리뷰 → merge'
    ])
    
    section_slide(prs, "10", "요약")
    content_slide(prs, "정리", [
        "git init / clone - 시작",
        "add + commit - 저장",
        "branch - 기능 개발",
        "push / pull - 협업",
        "이해가 안 되면 git status!"
    ])
    
    title_slide(prs, "다음 단계", "실전에서 연습하기")
    
    output_path = os.path.join(OUTPUT_DIR, "git-basics.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")

if __name__ == "__main__":
    create_presentation()
