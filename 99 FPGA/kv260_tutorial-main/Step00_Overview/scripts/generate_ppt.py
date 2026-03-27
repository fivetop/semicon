#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = "/home/woody/SoGang/2026-03/kv260_vivado_tutorial/Step00_Overview/docs"

PRIMARY = RGBColor(0, 51, 102)
ACCENT = RGBColor(255, 102, 0)
LIGHT = RGBColor(240, 248, 255)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(51, 51, 51)
GRAY = RGBColor(128, 128, 128)

def style_title(slide, title, subtitle=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = PRIMARY
    box.line.fill.background()
    
    tb = box.text_frame
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    if subtitle:
        p = tb.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(200, 200, 200)

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.8), Inches(13.333), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tb = sub_box.text_frame
    p = tb.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def section_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1.5))
    tb = num_box.text_frame
    p = tb.paragraphs[0]
    p.text = number
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(3), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

def content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(22)
        p.space_before = Pt(14)

def two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(5.8), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = PRIMARY
    left_card.line.width = Pt(2)
    
    ltb = left_card.text_frame
    p = ltb.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    for item in left_items:
        p = ltb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_before = Pt(10)
    
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5.8), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = ACCENT
    right_card.line.width = Pt(2)
    
    rtb = right_card.text_frame
    p = rtb.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    for item in right_items:
        p = rtb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_before = Pt(10)

def timeline_slide(prs, title, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    colors = [PRIMARY, RGBColor(0, 102, 153), ACCENT, RGBColor(0, 76, 153), PRIMARY]
    
    for i, step in enumerate(steps):
        x = Inches(0.6 + i * 2.6)
        y = Inches(2.2)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = colors[i]
        card.line.fill.background()
        
        num = slide.shapes.add_textbox(x + Inches(0.8), y + Inches(0.3), Inches(0.7), Inches(0.7))
        ntb = num.text_frame
        np = ntb.paragraphs[0]
        np.text = str(i + 1)
        np.font.size = Pt(24)
        np.font.bold = True
        np.font.color.rgb = WHITE
        np.alignment = PP_ALIGN.CENTER
        
        text = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.2), Inches(2), Inches(2))
        ttb = text.text_frame
        ttb.word_wrap = True
        tp = ttb.paragraphs[0]
        tp.text = step
        tp.font.size = Pt(16)
        tp.font.bold = True
        tp.font.color.rgb = WHITE
        tp.alignment = PP_ALIGN.CENTER

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    title_slide(prs, "KV260 Vivado 교육 과정", "Kria Vision AI Starter Kit으로 배우는 FPGA 설계")
    
    section_slide(prs, "01", "교육 개요")
    content_slide(prs, "커리큘럼 개요", [
        "KV260 FPGA 교육 프로그램 (5단계)",
        "대상: 반도체설계 신입 엔지니어",
        "일정: 2주 (10일)",
        "하드웨어: KV260 보드 (팀당 1대)"
    ])
    content_slide(prs, "교육 목표", [
        "Vivado 기반 FPGA 설계 역량 확보",
        "Verilog RTL 코딩 이해",
        "PS-PL 통신 구현",
        "비전 AI 플랫폼 이해"
    ])
    
    section_slide(prs, "02", "기술 스택")
    two_column_slide(prs, "개발 도구", "Vivado Design Suite 2021.2", [
        "Vivado Design Suite 2021.2",
        "Verilog HDL",
        "TCL 스크립트",
        "Python (PYNQ)"
    ], "하드웨어", [
        "KV260 Vision AI Starter Kit",
        "Zynq UltraScale+ MPSoC",
        "2GB DDR4",
        "MIPI CSI-2, DisplayPort"
    ])
    
    section_slide(prs, "03", "커리큘럼")
    timeline_slide(prs, "5단계 커리큘럼", [
        "LED Blink\n기초",
        "IP Integrator\n심화",
        "AXI GPIO\n통신",
        "Vitis HLS\n커스텀 IP",
        "BIST Platform\n고급"
    ])
    
    two_column_slide(prs, "Step 01-02 기초", "Step 01: LED Blink", [
        "Vivado 프로젝트 생성",
        "Verilog RTL 코딩",
        "블록 디자인 기초",
        "PYNQ에서 FPGA 프로그래밍"
    ], "Step 02: IP Integrator", [
        "Zynq UltraScale+ MPSoC",
        "Clocking Wizard",
        "플랫폼 익스포트 (XSA)",
        "Processor System Reset"
    ])
    
    two_column_slide(prs, "Step 03-04 중급", "Step 03: AXI GPIO", [
        "AXI GPIO IP 사용",
        "PS-PL 통신 구현",
        "Python으로 LED 제어",
        "인터럽트 기초"
    ], "Step 04: Vitis HLS", [
        "C/C++에서 RTL 생성",
        "커스텀 IP 개발",
        "PMOD 인터페이스",
        "PetaLinux 연동"
    ])
    
    content_slide(prs, "Step 05 - 고급", [
        "KV260 BIST Platform",
        "비전 AI 플랫폼 분석",
        "MIPI 카메라, VCU, DisplayPort",
        "PetaLinux 통합",
        "커스텀 애플리케이션 개발"
    ])
    
    section_slide(prs, "04", "준비물")
    two_column_slide(prs, "준비물", "하드웨어", [
        "KV260 Vision AI Starter Kit",
        "마이크로 USB 케이블",
        "이더넷 케이블",
        "SD 카드"
    ], "소프트웨어", [
        "Vivado 2021.2",
        "PYNQ 이미지",
        "VSCode + Verilog 확장",
        "Git"
    ])
    
    content_slide(prs, "학습 방법", [
        "이론 학습 (30분)",
        "실습 가이드 따라하기 (1시간)",
        "심화 연습 문제 (30분)",
        "검증 및 질문 (30분)"
    ])
    
    content_slide(prs, "기대 효과", [
        "Vivado로 FPGA 프로젝트 생성 가능",
        "Verilog로 RTL 설계 가능",
        "PS-PL 통신 구현 가능",
        "커스텀 IP 개발 가능",
        "비전 AI 애플리케이션 개발 기반 확보"
    ])
    
    section_slide(prs, "05", "다음 단계")
    title_slide(prs, "Step 01: LED Blink 기초", "Vivado로 첫 FPGA 프로젝트 만들기")
    
    output_path = os.path.join(OUTPUT_DIR, "course-intro.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")

if __name__ == "__main__":
    create_presentation()
