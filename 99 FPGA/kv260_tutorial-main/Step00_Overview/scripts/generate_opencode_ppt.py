#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = "/home/woody/SoGang/2026-03/kv260_vivado_tutorial/Step00_Overview/docs"

PRIMARY = RGBColor(0, 51, 102)
ACCENT = RGBColor(255, 102, 0)
LIGHT = RGBColor(240, 248, 255)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(128, 128, 128)
DARK = RGBColor(51, 51, 51)

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()
    
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.8), Inches(13.333), Inches(0.15))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tb = sub_box.text_frame
    p = tb.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def section_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()
    
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(1.5))
    tb = num_box.text_frame
    p = tb.paragraphs[0]
    p.text = number
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.2), Inches(3), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(2))
    tb = title_box.text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

def content_slide(prs, title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for item in items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(22)
        p.space_before = Pt(14)

def two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(5.8), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = WHITE
    left_card.line.color.rgb = PRIMARY
    left_card.line.width = Pt(2)
    
    ltb = left_card.text_frame
    p = ltb.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    for item in left_items:
        p = ltb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_before = Pt(10)
    
    right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5.8), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = WHITE
    right_card.line.color.rgb = ACCENT
    right_card.line.width = Pt(2)
    
    rtb = right_card.text_frame
    p = rtb.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    for item in right_items:
        p = rtb.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.space_before = Pt(10)

def code_slide(prs, title, code_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY
    
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(12.333), Inches(5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
    code_box.line.fill.background()
    
    tb = code_box.text_frame
    tb.word_wrap = True
    
    for line in code_lines:
        p = tb.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 255, 0)
        p.font.name = "Courier New"

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    title_slide(prs, "OpenCode & oh-my-opencode", "AI 코딩 도구 설치 및 활용")
    
    section_slide(prs, "01", "OpenCode란?")
    content_slide(prs, "OpenCode 이해", [
        "오픈소스 AI 코딩 에이전트",
        "터미널/데스크톱/IDE에서 사용 가능",
        "75+ LLM 프로바이더 지원 (Claude, GPT, Gemini 등)",
        "120K+ GitHub Stars, 5M+ 월간 사용자"
    ])
    
    two_column_slide(prs, "특징", "주요 기능", [
        "LSP 자동 로드",
        "멀티 세션 지원",
        "세션 공유 링크",
        "프로바이더 무관동일 사용"
    ], "장점", [
        "오픈소스 (무료)",
        "다양한 모델 지원",
        "프라이버시 우선",
        "커맨드라인 우선"
    ])
    
    section_slide(prs, "02", "OpenCode 설치")
    code_slide(prs, "설치 방법 (권장)", [
        '# curl 스크립트 (가장 빠름)',
        'curl -fsSL https://opencode.ai/install | bash',
        '',
        '# 또는 npm으로 설치',
        'npm install -g opencode-ai',
        '',
        '# Homebrew (macOS/Linux)',
        'brew install anomalyco/tap/opencode'
    ])
    
    two_column_slide(prs, "기타 설치 방법", "Node.js", [
        "npm install -g opencode-ai",
        "bun install -g opencode-ai",
        "pnpm install -g opencode-ai"
    ], "Go", [
        "go install github.com/opencode-ai/opencode@latest",
        "",
        "Arch Linux: paru -S opencode-bin",
        "Windows: choco install opencode"
    ])
    
    code_slide(prs, "설치 확인", [
        'opencode --version',
        'opencode --help',
        '',
        '## 출력 예시',
        '# OpenCode v0.x.x'
    ])
    
    section_slide(prs, "03", "API 키 설정")
    code_slide(prs, "환경 변수 설정", [
        '# Claude (Anthropic)',
        'export ANTHROPIC_API_KEY="your-key-here"',
        '',
        '# OpenAI (GPT)',
        'export OPENAI_API_KEY="your-key-here"',
        '',
        '# Google Gemini',
        'export GOOGLE_API_KEY="your-key-here"',
        '',
        '# ~/.bashrc 또는 ~/.zshrc에 추가하여 영구 설정'
    ])
    
    two_column_slide(prs, "설정 방법", "TUI로 설정", [
        "opencode 실행",
        "/connect 입력",
        "프로바이더 선택",
        "API 키 입력"
    ], "대화형 설정", [
        "opencode auth login",
        "Anthropic/OpenAI/Gemini 선택",
        "OAuth 또는 API 키 입력",
        "완료!"
    ])
    
    section_slide(prs, "04", "OpenCode 활용")
    code_slide(prs, "기본 명령어", [
        'opencode                    # 대화형 TUI 시작',
        'opencode run "코드 설명"  # 직접 실행',
        'opencode --prompt "..."    # 프롬프트',
        'opencode --model claude    # 모델 지정',
        'opencode models --refresh  # 모델 목록 새로고침'
    ])
    
    two_column_slide(prs, "주요 명령어", "세션 관리", [
        "/new - 새 세션 시작",
        "/continue - 이전 세션 계속",
        "/undo - 변경 취소",
        "/redo - 다시 실행"
    ], "모델 전환", [
        "/model claude-opus-4",
        "/model gpt-4o",
        "/model gemini-2.0-flash",
        "opencode --model ollama:codellama"
    ])
    
    content_slide(prs, "실용적인 활용", [
        "코드 설명: `opencode run \"Explain this function\"`",
        "버그 수정: `opencode run \"Fix this bug\"`",
        "리팩토링: `opencode run \"Refactor this code\"`",
        "테스트 생성: `opencode run \"Write tests for\"`"
    ])
    
    section_slide(prs, "05", "oh-my-opencode")
    content_slide(prs, "oh-my-opencode란?", [
        "OpenCode 기반 멀티 에이전트 오케스트레이션",
        "여러 AI 에이전트를 동시에 작업",
        "Plan → Explore → Build 모드",
        "프로젝트 전체 자동화"
    ])
    
    code_slide(prs, "설치 방법", [
        '# npm으로 전역 설치',
        'npm install -g oh-my-opencode',
        '',
        '# 또는 npx로 실행',
        'npx oh-my-opencode install',
        '',
        '# 옵션: TUI 없이 설치',
        'npx oh-my-opencode install --no-tui'
    ])
    
    two_column_slide(prs, "에이전트 시스템", "주요 에이전트", [
        "Sisyphus - 메인 오케스트레이터",
        "Prometheus - 플래닝/분석",
        "Oracle - 디버깅/아키텍처",
        "Momus - QA/검증"
    ], "실행 모드", [
        "Plan mode - 작업 계획",
        "Explore mode - 코드 분석",
        "Build mode - 실제 구현",
        "Analyze mode - 복합 분석"
    ])
    
    code_slide(prs, "기본 사용법", [
        '# 대시보드 열기',
        'oh-my-opencode',
        '',
        '# 특정 폴더에서 실행',
        'cd myproject && oh-my-opencode',
        '',
        '# 특정 작업 실행',
        'oh-my-opencode --task "Create a login feature"'
    ])
    
    section_slide(prs, "06", "설정 파일")
    code_slide(prs, "config.yaml 예시", [
        'model: claude-3-5-sonnet-20241022',
        '',
        'providers:',
        '  anthropic:',
        '    api_key: ${ANTHROPIC_API_KEY}',
        '  openai:',
        '    api_key: ${OPENAI_API_KEY}',
        '',
        'editor:',
        '  theme: terminal',
        '  tab_size: 2'
    ])
    
    content_slide(prs, "설정 파일 위치", [
        "~/.config/opencode/config.yaml",
        "프로젝트 디렉토리: .opencode/config.yaml",
        "환경 변수: ANTHROPIC_API_KEY, OPENAI_API_KEY"
    ])
    
    section_slide(prs, "07", "요약")
    content_slide(prs, "OpenCode 정리", [
        "설치: curl -fsSL https://opencode.ai/install | bash",
        "설정: opencode auth login 또는 환경 변수",
        "사용: opencode로 대화형 TUI 실행",
        "oh-my-opencode: 멀티 에이전트 자동화"
    ])
    
    two_column_slide(prs, "선택 가이드", "단순 코딩", [
        "OpenCode 직접 사용",
        "원하는 모델 선택 가능",
        "低成本"
    ], "복잡한 프로젝트", [
        "oh-my-opencode 사용",
        "멀티 에이전트 협업",
        "자동화된 플로우"
    ])
    
    title_slide(prs, "다음 단계", "실제로opencode 사용해보기")
    
    output_path = os.path.join(OUTPUT_DIR, "opencode-usage.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")

if __name__ == "__main__":
    create_presentation()
